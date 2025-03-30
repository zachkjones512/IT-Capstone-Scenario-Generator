'''
Disk File:
a. Password protected
b. Caesar cipher using {USER INPUT} shift left
c. base64 encoded password string
e. Find unique word in file
d. Steganography (USE PYTHON LIBRARY)
'''
from flask import Flask, request, send_file #using flask for webserver
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from caesarcipher import CaesarCipher
import pyzipper, os, random, base64

app = Flask(__name__) #creates flask application
pres = Presentation()
inputs = []
presName = 'test.pptx'
files = [5]
directoryName = "Scenario.zip"
randWords = [
    b"ZGV2ZWxvcGVy", b"Y29udGFpbmVy", b"YmlydGhkYXk=", b"cmFpbnN0b3Jt", 
    b"a2V5Ym9hcmQ=", b"cHJvZ3JhbQ==", b"bW91bnRhaW4=", b"YXN0cm9uYXV0",
    b"Z3VpdGFyaXN0", b"Y2hvY29sYXRl", b"YWR2ZW50dXJl", b"c25vd2ZhbGw=",
    b"Z2VuZXJhdGlvbg==", b"dW5pdmVyc2U=", b"ZGlzY292ZXI=", b"aW1hZ2luZQ==",
    b"bGFwdG9w", b"ZXZvbHV0aW9u", b"c3RhcmxpZ2h0", b"aG9zcGl0YWw=",
    b"ZGlhbW9uZA==", b"aW5ub3ZhdGlvbg==", b"d29ya2Zsb3c=", b"cGFpbnRpbmc="
]

def create_directory(inputslist):
    print("Creating directory...")
    pw = random.sample(randWords, 4)
    print(base64.b64decode(pw[3]).decode("utf-8"))

    with pyzipper.AESZipFile("level5.zip", 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as zipf: #FINISH
        zipf.setpassword(base64.b64decode(pw[3]).decode("utf-8").encode("utf-8"))
        with open("SENSITIVEINFO.txt", 'w') as file:
            file.write("DO NOT SHARE *** CONFIDENTIAL\nBelow are a list of passwords used by the company:\nYOU\nWIN!")
        zipf.write("SENSITIVEINFO.txt")
        print("Password set")


    with pyzipper.AESZipFile("level4.zip", 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as zipf: #WIP: SETGANOGRAPHY
        zipf.setpassword(base64.b64decode(pw[2]).decode("utf-8").encode("utf-8"))
        with open("password1.txt", 'w') as file:
            file.write("Test")
        zipf.write("password1.txt")
        zipf.write("level5.zip")
        os.remove("level5.zip")
        print("Password set")

    with pyzipper.AESZipFile("level3.zip", 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as zipf: #WIP: UNIQ 
        zipf.setpassword(base64.b64decode(pw[1]).decode("utf-8").encode("utf-8"))
        with open("password1.txt", 'w') as file:
            file.write("Test")
        zipf.write("password1.txt")
        zipf.write("level4.zip")
        os.remove("level4.zip")
        print("Password set")    

    with pyzipper.AESZipFile("level2.zip", 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as zipf: #BASE64
        zipf.setpassword(base64.b64decode(pw[0]).decode("utf-8").encode("utf-8"))
        with open("password1.txt", 'w') as file:
            file.write("VGhlIHBhc3N3b3JkIGlzOiA=" + pw[1].decode("utf-8"))
        zipf.write("password1.txt")
        zipf.write("level3.zip")
        os.remove("level3.zip")
        print("Password set")

    with pyzipper.AESZipFile("level1.zip", 'w', compression=pyzipper.ZIP_LZMA, encryption=pyzipper.WZ_AES) as zipf: #CAESAR
        zipf.setpassword(inputslist[1].encode("utf-8"))
        with open("password1.txt", 'w') as file:
            Cipher = CaesarCipher(base64.b64decode(pw[0]).decode("utf-8"),offset=int(inputs[3]))
            file.write("CAESAROffset=" + inputs[3] + "= '" + Cipher.encoded +"'")
        zipf.write("password1.txt")
        zipf.write("level2.zip")
        os.remove("level2.zip")
        print("Password set")

    with pyzipper.ZipFile(directoryName, 'w') as zipf: #THIS IS THE PARENT DIRECTORY
        zipf.write("level1.zip")
        zipf.write(presName)
        os.remove("password1.txt")
        os.remove("level1.zip")
        os.remove(presName)

def create_custom_slide(file_num,body_text): #creates custom slide layout used for every slide after title slide
    slide = pres.slides.add_slide(pres.slide_layouts[5])

    #create directory level number in top right    
    file_shape = slide.shapes.add_textbox(Inches(.1), Inches(.1), Inches(.5), Inches(.5))
    text_frame = file_shape.text_frame
    title_placeholder = text_frame.add_paragraph()
    title_placeholder.text = file_num 

    title_placeholder.alignment = PP_ALIGN.LEFT 
    title_placeholder.runs[0].font.size = Pt(12)
    title_placeholder.runs[0].font.color.rgb = RGBColor(255,255,255)

    #create body text
    body_shape = slide.shapes.add_textbox(Inches(.5), Inches(.75), Inches(9), Inches(6))
    body_frame = body_shape.text_frame
    body_frame.word_wrap = True 

    lines = body_text.split('\n') #allows for multiple paragraphs using newline operator"
    for line in lines: 
        body_placeholder = body_frame.add_paragraph()
        body_placeholder.text = line
        body_placeholder.alignment = PP_ALIGN.CENTER 
        body_placeholder.runs[0].font.size = Pt(20)
        body_placeholder.runs[0].font.color.rgb = RGBColor(0, 143, 17)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,0,0)

    return slide

@app.route('/submission', methods=['POST']) #user clicks submit on any of the questions
def submission():
    global inputs
    usrStr = request.get_json() #submission button sends post request with answers json file
    inputs = usrStr.get('responses', [])
    return "Input received, " + inputs[0]

@app.route('/complete', methods=['GET']) #user completes the questionnaire
def complete():
    if not inputs:
        return "No inputs received. Please submit some data first.", 400 

    body_text_dict = [
        "Security breach detected! Someone has stolen sensitive information from CyberTech Manufacturing in " + inputs[5] + \
        "\n \n Police have narrowed down the search to one suspect: " + inputs[1] + ". \n \nFollowing a warrant they have seized a USB drive with an encrypted file directory." \
        " They need your help in retrieving the lost data! \n \n" + inputs[1] + "'s colleagues say that they followed bad cybersecurity practices and often used their own name as the password to their computer",
    #
        inputs[2],
        inputs[3],
        inputs[4],
        inputs[5],
        inputs[6],    
    ]

    #Create title slide
    create_custom_slide("Here is " + inputs[0] +"'s Cybersecurity Scenario!",
                                 "For this cybersecurity scenario you will be taken through an investigation into the file directory that you downloaded with this PowerPoint. You are tasked with navigating each level by finding the subsequent directory's password." \
    "\n \n \nThere are a couple requirements before you begin. First, you need a file explorer capable of opening password-protected directories. " \
    "We recommend the free and open source 7-zip! \n \nThe second requirement is a linux terminal, "\
    "macOS machines contain this already, while Windows machines will need to use a virtual machine or Windows Subsystem for Linux (WSL)\n \nWhen you are ready to begin, go to the next slide!")
    print("Created presentation")
    #Title slide created

    for i in range(1,7):
        slide = create_custom_slide("Level " + str(i), body_text_dict[i-1])
        
        hint_shape = slide.shapes.add_textbox(Inches(.1), Inches(6.5), Inches(4), Inches(.5))
        hint_frame = hint_shape.text_frame
        hint_text = hint_frame.add_paragraph()
        hint_text.text = "Having trouble? Open the file hint.txt in this directory level!"

        hint_text.alignment = PP_ALIGN.LEFT #aligns file directory height to top left
        hint_text.runs[0].font.size = Pt(12)
        hint_text.runs[0].font.color.rgb = RGBColor(255,0,0)

    try:
        pres.save(presName)
        print("Saved successfully")
    except:
        return "Error saving presentation", 500

    create_directory(inputs)
    inputs.clear()

    try:
        return send_file("../" + directoryName, mimetype='zip', as_attachment=True, download_name=directoryName)
    except:
        return "Error uploading file. Please try again later."




if __name__ == '__main__':
    app.run(port=5500) #listens on port 5500
